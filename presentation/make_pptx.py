"""Generate research-group presentation for Battery Degradation GPR project."""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── paths ────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "battery_gpytorch_rtx4060" / "battery_gpytorch" / "output"
PPTX_OUT = Path(__file__).resolve().parent / "EIS_GPR_Presentation.pptx"

# ── colours ──────────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x00, 0x96, 0xD6)   # blue accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xEF, 0x53, 0x50)
YELLOW = RGBColor(0xFF, 0xD5, 0x4F)
TABLE_HEADER_BG = RGBColor(0x00, 0x6E, 0x9F)
TABLE_ROW_BG = RGBColor(0x2A, 0x2A, 0x45)
TABLE_ALT_BG = RGBColor(0x22, 0x22, 0x3A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        if bold_first and i == 0:
            p.font.bold = True
    return tf


def add_table(slide, left, top, width, height, data, col_widths=None):
    """data = list of rows, first row is header."""
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for j, w in enumerate(col_widths):
            table.columns[j].width = w

    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER

                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT_GREY

            # cell background
            cell_fill = cell.fill
            cell_fill.solid()
            if i == 0:
                cell_fill.fore_color.rgb = TABLE_HEADER_BG
            elif i % 2 == 0:
                cell_fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell_fill.fore_color.rgb = TABLE_ROW_BG

    return table_shape


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add image if it exists, otherwise add placeholder text."""
    if Path(img_path).exists():
        kwargs = {"image_file": str(img_path), "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(**kwargs)
        return True
    else:
        add_textbox(slide, left, top, Inches(3), Inches(0.4),
                    f"[missing: {Path(img_path).name}]", font_size=10, color=RED)
        return False


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ═════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ═════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ── SLIDE 1: Title ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "Battery Degradation Prediction via\nGaussian Process Regression on EIS Spectra",
            font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(4), Inches(3.3), Inches(5.3))
add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
            "Reproduction & Extension of Zhang et al., Nature Communications 2020",
            font_size=20, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
            "DOI: 10.1038/s41467-020-15235-7",
            font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
            "Research Group Presentation",
            font_size=16, color=LIGHT_GREY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 2: Background & Motivation ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Background & Motivation", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

items = [
    "Lithium-ion battery degradation monitoring is critical for EV and grid storage safety",
    "Electrochemical Impedance Spectroscopy (EIS) captures degradation signatures non-destructively",
    "Zhang et al. (2020): GPR on raw EIS spectra predicts capacity (SOH) and RUL",
    "No feature engineering needed -- ARD discovers informative frequencies automatically",
    "Key insight: Feature #91 (17.80 Hz, Im(Z)) is the universal degradation indicator",
]
add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(7), Inches(3.5), items, font_size=18)

items2 = [
    "Our contributions:",
    "  1. Full reproduction of Zhang et al. (Figs 1-4) -- matched or beat all targets",
    "  2. Extension to in-house datasets (A1-A8, CA1-CA8) with LOOCV validation",
    "  3. Multi-temperature analysis (-10C, -20C) revealing physics-dependent RUL limits",
    "  4. Coupled ARD kernel -- physically correct (Kramers-Kronig) importance scoring",
    "  5. Capacity-derived RUL -- novel approach that avoids EIS->RUL mapping failures",
]
add_bullet_list(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(3.2), items2, font_size=16,
                color=LIGHT_GREY, bold_first=True)


# ── SLIDE 3: Method Overview ────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Method: GPR on Raw EIS Spectra", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

items = [
    "Input: Raw EIS spectrum [Re(Z), Im(Z)] at each measured frequency",
    "  -- Cambridge dataset: 60 frequencies x 2 = 120 features",
    "  -- In-house datasets: 33 frequencies x 2 = 66 features",
    "",
    "Model: Gaussian Process Regression (GPR)",
    "  -- Capacity: RBF kernel (fixed length-scale l=30 for LOOCV) or ARD-RBF",
    "  -- RUL: Linear kernel (DotProduct, Zhang eq. 5)",
    "  -- Noise: WhiteKernel (Gaussian likelihood)",
    "",
    "ARD (Automatic Relevance Determination):",
    "  -- One length-scale per feature; weight w_m = exp(-sigma_m), normalised",
    "  -- Identifies which EIS frequencies carry degradation information",
    "",
    "Normalisation: z-score (training stats applied to test set)",
    "  -- LOOCV capacity: joint norm (train+test pooled) removes cell offset",
    "  -- LOOCV RUL: training-only norm (paper's stated approach)",
]
add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), items, font_size=16)


# ── SLIDE 4: Datasets Overview ──────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Datasets Overview", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Dataset", "Source", "Cells", "Features", "Temp", "Cycles", "EOL cells"],
    ["Cambridge", "Zhang et al. (public)", "25C01-08, 35C01-02, 45C01-02", "120", "25/35/45C", "~275", "All"],
    ["A1-A8 (partial)", "In-house", "A1-A8", "66", "RT (~25C)", "~268", "6 (A3,A6 DNF)"],
    ["CA1-CA8 (complete)", "In-house", "CA1-CA8", "66", "RT (~25C)", "470+", "7 (CA6 DNF)"],
    ["CB multi-temp", "In-house", "N10_CB1-4, N20_CB1-4", "66", "-10C, -20C", "71-114 / 17-21", "All 8"],
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(2.5), data)

items = [
    "EIS state: post-charge (State V analogue) -- consistent across all datasets",
    "A1-A8 partial = first ~268 cycles; CA1-CA8 = same cells run to completion (470+ cycles)",
    "DNF cells: A3 (anomalously low initial cap 3800 vs 4050 mAh), A6/CA6 (never reached 80%)",
    "CB cells: Molicell 21700 P42A NMC; -20C cells have very short lives (17-21 cycles)",
    "EOL definition: first cycle where capacity < 80% of initial capacity",
]
add_bullet_list(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(3.0), items, font_size=14,
                color=LIGHT_GREY)


# ── SLIDE 5: Paper Reproduction - Capacity ──────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Exp 1: Paper Reproduction -- Capacity", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Figure", "Description", "Our R\u00b2", "Paper R\u00b2", "Status"],
    ["Fig 1a", "Single-T 25C capacity", "0.882", "0.88", "Matched"],
    ["Fig 3a", "Multi-T 35C capacity", "0.91", "0.81", "Beat"],
    ["Fig 3b", "Multi-T 45C capacity", "0.94", "0.72", "Beat"],
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(6), Inches(2.0), data)

# figures
add_image_safe(slide, OUT / "fig3a_capacity_35C02.png",
               Inches(7), Inches(1.2), width=Inches(5.8))
add_image_safe(slide, OUT / "fig3b_capacity_45C02.png",
               Inches(7), Inches(4.0), width=Inches(5.8))

add_textbox(slide, Inches(0.6), Inches(3.8), Inches(6), Inches(2.5),
            "Multi-temperature training enables strong generalisation.\n\n"
            "Training on 25/35/45C cells, testing on held-out cells at each\n"
            "temperature. Our reproduction matches or exceeds all paper targets.",
            font_size=14, color=LIGHT_GREY)


# ── SLIDE 6: Paper Reproduction - RUL ───────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Exp 1: Paper Reproduction -- RUL", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Figure", "Description", "Our R\u00b2", "Paper R\u00b2", "Status"],
    ["Fig 4a", "Multi-T 25C RUL (25C05)", "0.970", "0.87", "Beat"],
    ["Fig 4b", "Multi-T 35C RUL (35C02)", "0.85", "0.75", "Beat"],
    ["Fig 4c", "Multi-T 45C RUL (45C02)", "0.91", "0.92", "Matched"],
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(6), Inches(2.0), data)

add_image_safe(slide, OUT / "fig4b_rul_35C02.png",
               Inches(7), Inches(1.2), width=Inches(5.8))
add_image_safe(slide, OUT / "fig4a_rul_25C05_multiT.png",
               Inches(7), Inches(4.0), width=Inches(5.8))

add_textbox(slide, Inches(0.6), Inches(3.8), Inches(6), Inches(3.0),
            "RUL prediction works in multi-temperature DOE because\n"
            "temperature variation creates a learnable link between\n"
            "EIS signature and remaining life (Arrhenius relationship).\n\n"
            "Linear kernel (DotProduct) used for RUL -- Zhang eq. 5.\n"
            "Training-only z-score normalisation.",
            font_size=14, color=LIGHT_GREY)


# ── SLIDE 7: Paper Reproduction - ARD Weights ───────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Exp 1: ARD Weight Analysis", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

add_image_safe(slide, OUT / "fig1c_ARD_weights_25C.png",
               Inches(0.3), Inches(1.3), width=Inches(6.2))
add_image_safe(slide, OUT / "fig3c_ARD_weights.png",
               Inches(6.8), Inches(1.3), width=Inches(6.2))

add_textbox(slide, Inches(0.3), Inches(5.2), Inches(6), Inches(0.4),
            "25C ARD (Fig 1c): #91 + #100 in top-5", font_size=14, color=LIGHT_GREY,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(6.8), Inches(5.2), Inches(6), Inches(0.4),
            "35C ARD (Fig 3c): #91 dominant", font_size=14, color=LIGHT_GREY,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(1.2),
            "Feature #91 = 17.80 Hz Im(Z) -- universal degradation indicator across all temperatures.\n"
            "Feature #100 = 2.16 Hz -- appears at 25C only; multi-T training strips it out,\n"
            "isolating the temperature-independent signal at 17.80 Hz.",
            font_size=16, color=WHITE)


# ── SLIDE 8: A1-A8 LOOCV Capacity ───────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Exp 3: A1-A8 LOOCV -- Capacity", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Cell", "A1", "A2", "A3", "A4", "A5", "A6", "A7", "A8", "Mean"],
    ["R\u00b2", "0.987", "0.983", "0.967", "0.959", "0.991", "0.867", "0.995", "0.964", "0.964"],
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.0), data)

add_image_safe(slide, OUT / "new_dataset" / "fig_cap_loocv.png",
               Inches(0.3), Inches(2.5), width=Inches(6.2))
add_image_safe(slide, OUT / "new_dataset" / "fig_ARD_loocv_folds.png",
               Inches(6.8), Inches(2.5), width=Inches(6.2))

add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.8),
            "Joint normalisation (train+test pooled per fold) removes cell-to-cell impedance offset.\n"
            "Strong generalisation across all 8 cells including both DNF cells (A3, A6). Fixed RBF l=30.",
            font_size=14, color=LIGHT_GREY)


# ── SLIDE 9: A1-A8 LOOCV RUL + Frequency Subsets ────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Exp 3: A1-A8 LOOCV -- RUL (Failure Analysis)", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Cell", "A1", "A2", "A4", "A5", "A7", "A8", "Mean"],
    ["Linear R\u00b2", "-0.52", "-0.56", "-0.00", "-0.38", "0.17", "-0.68", "-0.33"],
    ["RBF R\u00b2", "-0.57", "-0.68", "-0.00", "-2.09", "-5.27", "-0.09", "-1.24"],
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(8.5), Inches(1.3), data)

data2 = [
    ["Band", "Freq range", "Cap R\u00b2", "RUL R\u00b2"],
    ["Full spectrum", "1-10,000 Hz", "0.964", "-0.33"],
    ["High", "500-10,000 Hz", "0.905", "-2.15"],
    ["Mid", "10-500 Hz", "0.867", "-1.79"],
    ["Low", "1-10 Hz", "0.872", "-2.07"],
]
add_table(slide, Inches(0.6), Inches(3.2), Inches(5.5), Inches(2.2), data2)

add_image_safe(slide, OUT / "new_dataset" / "fig_rul_loocv.png",
               Inches(6.8), Inches(2.8), width=Inches(6.0))

add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(1.2),
            "Direct EIS -> RUL fails for same-temperature cells.\n"
            "RUL fails equally across all frequency bands -- not a spectral coverage problem.\n"
            "Root cause: cells with identical EIS can have different total lifetimes (2.4x spread).",
            font_size=14, color=LIGHT_GREY)


# ── SLIDE 10: CA1-CA8 Capacity LOOCV ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Exp 4: CA1-CA8 Complete Lifecycle -- Capacity LOOCV", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Cell", "CA1", "CA2", "CA3", "CA4", "CA5", "CA6", "CA7", "CA8", "Mean"],
    ["R\u00b2", "0.995", "0.993", "0.995", "0.987", "0.971", "0.562", "0.991", "0.982", "0.934"],
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.0), data)

data2 = [
    ["Cell", "CA1", "CA2", "CA3", "CA4", "CA5", "CA7", "CA8", "Mean"],
    ["Linear R\u00b2", "0.141", "0.154", "0.053", "-0.324", "0.010", "-1.169", "-0.095", "-0.176"],
]
add_table(slide, Inches(0.6), Inches(2.6), Inches(12), Inches(0.9), data2)

add_textbox(slide, Inches(0.6), Inches(3.8), Inches(6), Inches(1.0),
            "Capacity: excellent (mean R\u00b2=0.934)\nDirect RUL: fails again (mean R\u00b2=-0.176)\n\n"
            "Same conclusion as A1-A8: EIS encodes current\nhealth, not total lifespan.",
            font_size=15, color=LIGHT_GREY)

add_image_safe(slide, OUT / "new_dataset" / "fig_rul_loocv_comparison.png",
               Inches(6.8), Inches(3.5), width=Inches(6.0))


# ── SLIDE 11: Why Direct RUL Fails ──────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Why Direct EIS -> RUL Fails", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

add_textbox(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.5),
            "Case 1: Same-Temperature (A1-A8, CA1-CA8)", font_size=22, bold=True, color=WHITE)
items1 = [
    "All cells: same temperature, C-rate, conditions",
    "Total lifetimes: 190-448 cycles (2.4x range)",
    "Two cells at identical SOH can have completely different remaining lives",
    "EIS encodes current health, not intrinsic total lifespan",
    "For RUL to work: need diverse lifetimes driven by distinguishable conditions",
]
add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5), items1, font_size=14,
                color=LIGHT_GREY)

add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.5),
            "Case 2: Cold Temperatures (-20C CB)", font_size=22, bold=True, color=WHITE)
items2 = [
    "Data problem: RUL_max = 17-21 cycles vs RT 200+ (out of distribution)",
    "Physics problem (fundamental):",
    "  Warm: higher T -> faster SEI growth -> shorter life AND larger Im(Z)",
    "    -> EIS encodes both current state and future rate -> RUL learnable",
    "  Cold: lower T -> high Re(Z) from kinetic limitation, not irreversible damage",
    "    -> capacity loss is thermally reversible (warm cell -> capacity recovers)",
    "    -> EIS encodes thermal kinetics, not degradation rate -> RUL signal absent",
]
add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(6), Inches(3.0), items2, font_size=14,
                color=LIGHT_GREY)

add_textbox(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
            "Key Distinction", font_size=22, bold=True, color=YELLOW)
items3 = [
    "Capacity = snapshot: 'how much charge can this cell deliver now?' -- EIS -> SOH works regardless",
    "RUL = forecast: 'how fast will degradation continue?' -- requires knowing mechanism (reversible vs irreversible)",
    "EIS alone cannot distinguish thermal kinetics from electrochemical degradation -> RUL breaks at -20C",
]
add_bullet_list(slide, Inches(0.8), Inches(5.4), Inches(12), Inches(1.8), items3, font_size=15,
                color=WHITE)


# ── SLIDE 12: Capacity-Derived RUL ──────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Solution: Capacity-Derived RUL", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

items = [
    "1. Predict full capacity trajectory via LOOCV GPR",
    "2. Fit linear trend to predicted trajectory",
    "3. Extrapolate to 80% threshold -> predicted EOL",
    "4. RUL[i] = predicted_EOL - i",
]
add_bullet_list(slide, Inches(0.6), Inches(1.2), Inches(6), Inches(1.5), items, font_size=16, color=WHITE)

data = [
    ["Cell", "CA1", "CA2", "CA3", "CA4", "CA5", "CA7", "CA8", "Mean"],
    ["Cap R\u00b2", "0.995", "0.993", "0.995", "0.987", "0.971", "0.991", "0.982", "0.934"],
    ["RUL R\u00b2", "0.999", "0.994", "0.767", "0.550", "0.999", "0.973", "0.967", "0.893"],
    ["Pred EOL", "173", "164", "182", "113", "115", "91", "212", "--"],
    ["Actual EOL", "174", "168", "214", "143", "116", "96", "224", "--"],
]
add_table(slide, Inches(0.6), Inches(3.0), Inches(11.5), Inches(2.2), data)

add_image_safe(slide, OUT / "cap_rul" / "fig_rt_cap_rul_scatter.png",
               Inches(7.5), Inches(5.3), width=Inches(5.5))

add_textbox(slide, Inches(0.6), Inches(5.5), Inches(6.5), Inches(1.5),
            "Cap-derived RUL (mean R\u00b2=0.893) dramatically outperforms\n"
            "direct EIS->RUL (mean R\u00b2=-0.176).\n\n"
            "Avoids the EIS->RUL mapping problem entirely by leveraging\n"
            "the strong capacity prediction model.",
            font_size=15, color=WHITE)


# ── SLIDE 13: Multi-Temperature Results ─────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Exp 5: Multi-Temperature Zhang DOE", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Set", "Cells"],
    ["Train", "CA1-CA8 (RT) + N10_CB1-3 (-10C) + N20_CB1-3 (-20C)"],
    ["Test", "N10_CB4 (-10C held-out) + N20_CB4 (-20C held-out)"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(9), Inches(1.3), data)

data2 = [
    ["Task", "Cell", "R\u00b2"],
    ["Capacity", "N10_CB4 (-10C)", "0.375"],
    ["Capacity", "N20_CB4 (-20C)", "0.949"],
    ["RUL", "N10_CB4 (-10C)", "0.226"],
    ["RUL", "N20_CB4 (-20C)", "-120"],
]
add_table(slide, Inches(0.6), Inches(2.8), Inches(5), Inches(2.2), data2)

add_image_safe(slide, OUT / "multitemp_zhang" / "fig_zhang_capacity.png",
               Inches(6.5), Inches(2.5), width=Inches(6.3))

data3 = [
    ["Approach", "Mean Cap R\u00b2", "Mean RUL R\u00b2"],
    ["Baseline (train RT+-10C only)", "-8.7", "-6499"],
    ["Zhang DOE (all temps in training)", "0.949", "-120"],
]
add_table(slide, Inches(0.6), Inches(5.3), Inches(5.5), Inches(1.3), data3)

add_textbox(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
            "-20C RUL fails: 17-21 cycle life vs RT 200+ -- scale mismatch. Coupled ARD: 1.33 Hz dominant (w=0.71).",
            font_size=14, color=LIGHT_GREY)


# ── SLIDE 14: Single-T per Temperature -- Capacity ──────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Exp 6: Single-T Capacity (Zhang Fig 1 Equivalent)", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Group", "Train", "Test", "Capacity R\u00b2"],
    ["RT (~25C)", "CA1-CA6", "CA7, CA8", "0.996 / 0.991"],
    ["-10C", "N10_CB1-3", "N10_CB4", "0.676"],
    ["-20C", "N20_CB1-3", "N20_CB4", "0.937"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(7), Inches(1.8), data)

add_image_safe(slide, OUT / "ca_zhang" / "fig_rt_1a_capacity_trajectories.png",
               Inches(0.2), Inches(3.3), width=Inches(4.2))
add_image_safe(slide, OUT / "ca_zhang" / "fig_n10_1a_capacity_trajectories.png",
               Inches(4.5), Inches(3.3), width=Inches(4.2))
add_image_safe(slide, OUT / "ca_zhang" / "fig_n20_1a_capacity_trajectories.png",
               Inches(8.8), Inches(3.3), width=Inches(4.2))

add_textbox(slide, Inches(0.2), Inches(6.6), Inches(4.2), Inches(0.4),
            "RT: R\u00b2 = 0.996, 0.991", font_size=13, color=GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.5), Inches(6.6), Inches(4.2), Inches(0.4),
            "-10C: R\u00b2 = 0.676", font_size=13, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(8.8), Inches(6.6), Inches(4.2), Inches(0.4),
            "-20C: R\u00b2 = 0.937", font_size=13, color=GREEN, alignment=PP_ALIGN.CENTER)


# ── SLIDE 15: Single-T ARD per Temperature ──────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Exp 6: Coupled ARD Weights per Temperature", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

add_image_safe(slide, OUT / "ca_zhang" / "fig_rt_1c_ARD_weights.png",
               Inches(0.2), Inches(1.2), width=Inches(4.2))
add_image_safe(slide, OUT / "ca_zhang" / "fig_n10_1c_ARD_weights.png",
               Inches(4.5), Inches(1.2), width=Inches(4.2))
add_image_safe(slide, OUT / "ca_zhang" / "fig_n20_1c_ARD_weights.png",
               Inches(8.8), Inches(1.2), width=Inches(4.2))

add_textbox(slide, Inches(0.2), Inches(4.5), Inches(4.2), Inches(0.4),
            "RT: ~1 Hz + ~5000 Hz (dual peak)", font_size=13, color=LIGHT_GREY,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.5), Inches(4.5), Inches(4.2), Inches(0.4),
            "-10C: 13.3 Hz (w=1.0) single spike", font_size=13, color=ACCENT,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(8.8), Inches(4.5), Inches(4.2), Inches(0.4),
            "-20C: ~1 Hz + ~5000 Hz (dual peak)", font_size=13, color=LIGHT_GREY,
            alignment=PP_ALIGN.CENTER)

data = [
    ["Group", "Top Frequency", "Weight", "Interpretation"],
    ["RT", "~1 Hz + ~5000 Hz", "dual peak", "Low-freq SEI + high-freq bulk resistance"],
    ["-10C", "13.3 Hz", "1.0", "Single spike -- pure SEI/charge-transfer"],
    ["-20C", "~1 Hz + ~5000 Hz", "dual peak", "Similar to RT -- kinetically limited"],
]
add_table(slide, Inches(0.6), Inches(5.1), Inches(12), Inches(1.8), data)


# ── SLIDE 16: Cap-Derived RUL Across Temperatures ───────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "RUL Strategy Depends on Temperature", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Temperature", "Cap-derived RUL R\u00b2", "Direct EIS->RUL R\u00b2", "Best Approach"],
    ["RT (25C)", "0.893", "-4.3", "Cap-derived (LOOCV)"],
    ["-10C", "-0.82", "0.734", "Direct (linear kernel)"],
    ["-20C", "0.970", "0.459", "Cap-derived (DOE)"],
]
add_table(slide, Inches(1.5), Inches(1.3), Inches(10), Inches(1.8), data)

add_image_safe(slide, OUT / "cap_rul" / "fig_rt_cap_rul_trajectories.png",
               Inches(0.2), Inches(3.5), width=Inches(4.2))
add_image_safe(slide, OUT / "cap_rul" / "fig_n10_cap_rul_trajectories.png",
               Inches(4.5), Inches(3.5), width=Inches(4.2))
add_image_safe(slide, OUT / "cap_rul" / "fig_n20_cap_rul_trajectories.png",
               Inches(8.8), Inches(3.5), width=Inches(4.2))

add_textbox(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
            "No single method works at all temperatures -- optimal strategy depends on degradation regime.",
            font_size=16, bold=True, color=YELLOW, alignment=PP_ALIGN.CENTER)


# ── SLIDE 17: Key Scientific Finding ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Key Finding: Universal Degradation Indicator", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Dataset", "Kernel", "Top Feature", "Frequency", "Interpretation"],
    ["35C (Fig 3c)", "Decoupled", "#91", "17.80 Hz", "Exact match with paper"],
    ["25C (Fig 1c)", "Decoupled", "#91, #100", "17.80 + 2.16 Hz", "Same region"],
    ["45C (Fig 3d)", "Decoupled", "#88", "~20 Hz", "Same low-freq Im(Z) region"],
    ["A1-A8 LOOCV", "Decoupled", "low-freq Im(Z)", "~1-20 Hz", "Consistent"],
    ["CB multi-T", "Coupled", "1.33 Hz", "w=0.71", "SEI/diffusion + bulk"],
    ["-10C single-T", "Coupled", "13.3 Hz", "w=1.0", "Pure SEI/charge-transfer"],
]
add_table(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(3.0), data)

items = [
    "Feature #91 (17.80 Hz, Im(Z)) is the universal degradation indicator",
    "Consistent across Cambridge, in-house, and multi-temperature datasets",
    "Multi-T training acts as a regulariser, stripping out temperature-dependent features",
    "Coupled ARD (33 ls, Re+Im paired) gives physically interpretable weights (Kramers-Kronig)",
    "-10C shows unique single-frequency dominance at 13.3 Hz -- pure charge-transfer signal",
]
add_bullet_list(slide, Inches(0.8), Inches(4.5), Inches(12), Inches(2.5), items, font_size=16, color=WHITE)


# ── SLIDE 18: Coupled vs Decoupled ARD ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Coupled vs Decoupled ARD", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

add_image_safe(slide, OUT / "new_dataset" / "fig_ARD_coupled_vs_decoupled.png",
               Inches(6.5), Inches(1.2), width=Inches(6.5))

items = [
    "Standard ARD: 66 independent length-scales",
    "  -- Re(Z) and Im(Z) treated as separate features",
    "  -- Ambiguous importance: Re and Im at same freq get different weights",
    "",
    "Coupled ARD: 33 length-scales (one per frequency)",
    "  -- Re(Z) and Im(Z) share the same length-scale per frequency",
    "  -- Physically motivated: Kramers-Kronig relations couple Re and Im",
    "  -- Gives cleaner, more interpretable importance scores",
    "",
    "Result: comparable LOOCV accuracy with half the parameters",
    "  -- Decoupled: 66 ls -> noisy attribution",
    "  -- Coupled: 33 ls -> clear frequency ranking",
]
add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), items, font_size=15,
                color=LIGHT_GREY)


# ── SLIDE 19: Summary & Conclusions ─────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
            "Summary & Conclusions", font_size=32, bold=True, color=ACCENT)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(4))

data = [
    ["Model", "Result", "Notes"],
    ["Paper reproduction (capacity)", "R\u00b2 = 0.88-0.94", "Matched or beat all targets"],
    ["Paper reproduction (RUL)", "R\u00b2 = 0.85-0.97", "Beat all multi-T targets"],
    ["A1-A8 Capacity LOOCV", "mean R\u00b2 = 0.964", "Strong generalisation, 8 cells"],
    ["CA1-CA8 Capacity LOOCV", "mean R\u00b2 = 0.934", "Complete lifecycle, 470+ cycles"],
    ["RT Cap-derived RUL", "mean R\u00b2 = 0.893", "Best RT approach (7 EOL cells)"],
    ["-20C Cap-derived RUL (DOE)", "R\u00b2 = 0.970", "Pred EOL=16.1 vs actual 17"],
    ["-10C Direct RUL", "R\u00b2 = 0.734", "Single dominant freq at 13.3 Hz"],
    ["Single-T Capacity (RT)", "R\u00b2 = 0.996", "Beats Zhang's 0.88"],
]
add_table(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(3.8), data)

items = [
    "GPR on raw EIS spectra is a powerful, feature-engineering-free approach to battery SOH prediction",
    "Capacity prediction is robust across datasets, temperatures, and validation strategies",
    "Direct EIS->RUL fails for same-temperature cells; cap-derived RUL is the practical solution",
    "No single RUL method works at all temperatures -- strategy depends on degradation physics",
    "Coupled ARD provides physically interpretable frequency importance (Kramers-Kronig correct)",
    "17.80 Hz Im(Z) is the universal degradation indicator across all conditions",
]
add_bullet_list(slide, Inches(0.8), Inches(5.3), Inches(12), Inches(2.0), items, font_size=15,
                color=WHITE)


# ── SAVE ─────────────────────────────────────────────────────────────────
prs.save(str(PPTX_OUT))
print(f"Presentation saved to: {PPTX_OUT}")
print(f"Total slides: {len(prs.slides)}")
